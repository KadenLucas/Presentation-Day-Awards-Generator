from pptx import Presentation
from io import BytesIO, StringIO
from zipfile import ZipFile
import csv


def process_awards(awards_file: BytesIO) -> dict[str, list]:
    award_categories = dict()
    header = None

    for line in csv.reader(StringIO(awards_file.read().decode())):
        if header is None:
            header = line
            continue
        
        entry = {
            field: value for field, value in zip(header, line)
        }

        award_type = entry["Award Type"]
        awards = award_categories.get(award_type)

        if awards is None:
            awards = list()
            award_categories[award_type] = awards
        
        awards.append(entry)
    
    return award_categories


def generate(template_fp: str, calendar_year: int, awards: BytesIO):
    categories = process_awards(awards)
    zip_buffer = BytesIO()

    with ZipFile(zip_buffer, 'w') as zip:
        for category, awards in categories.items():
            template = Presentation(template_fp)
            layout_name = category

            for award in sorted(awards, key=lambda x: x["Surname"]):
                award_slide = template.slides.add_slide(
                    template.slide_layouts.get_by_name(layout_name)
                )

                name, academic_year, subject, calendar_year = award_slide.placeholders
                name.text = f"{award['First Name']} {award['Surname']}"
                academic_year.text = f"Year {award['Year']}"
                subject.text = award["Subject"]
                calendar_year.text = "2025"

            output = BytesIO()
            template.save(output)
            output.seek(0)
            zip.writestr(f"{category}.pptx", output.read())

    zip_buffer.seek(0)
    return zip_buffer
    
if __name__ == "__main__":
    with open("assets/sample.csv", 'r') as awards:
        generate(
            "assets/template.pptx",
            2025, awards
        )
